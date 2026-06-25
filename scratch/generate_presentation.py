import sys
import subprocess
import os

# 1. Dynamically install python-pptx if it is not present
try:
    import pptx
except ImportError:
    print("python-pptx not found. Installing dynamically...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # Initialize Presentation
    prs = Presentation()
    
    # Set to 16:9 widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 2. Design System Colors
    COLOR_BG = RGBColor(248, 250, 252)       # Light Slate Bg
    COLOR_PRIMARY = RGBColor(5, 150, 105)     # Emerald Green
    COLOR_TEXT_DARK = RGBColor(30, 41, 59)   # Slate 800 (Dark text)
    COLOR_TEXT_LIGHT = RGBColor(100, 116, 139) # Slate 500 (Subtext)
    COLOR_ACCENT = RGBColor(217, 119, 6)     # Amber Accent
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_HEADER_BG = RGBColor(15, 23, 42)    # Dark Slate 900
    
    # Helper to apply clean background to slides
    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper to add standard Header banner
    def add_slide_header(slide, title_text, category_text="CARBONINTEL PROJECT REVIEW"):
        # Header Dark Banner
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = COLOR_HEADER_BG
        header_shape.line.fill.background()
        
        # Category label (e.g. "CARBONINTEL PROJECT REVIEW")
        cat_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.12), Inches(12), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = 'Arial'
        p_cat.font.size = Pt(9)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_PRIMARY
        
        # Slide Title text
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.6))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = 'Arial'
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_WHITE

    # Helper to add Footer
    def add_slide_footer(slide, slide_num):
        footer_box = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(12.133), Inches(0.4))
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"CarbonIntel: Precision Agriculture Carbon Intelligence Platform    |    Slide {slide_num}"
        p.font.name = 'Arial'
        p.font.size = Pt(9)
        p.font.color.rgb = COLOR_TEXT_LIGHT
        p.alignment = PP_ALIGN.LEFT

    # ==========================================
    # SLIDE 1: Title Slide (Dark Theme)
    # ==========================================
    blank_slide_layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1, COLOR_HEADER_BG)
    
    # Left decorative accent line
    accent_bar = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(0.15), Inches(3.2)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_PRIMARY
    accent_bar.line.fill.background()
    
    # Title Box
    title_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(11.0), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "CarbonIntel"
    p.font.name = 'Arial'
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    p2 = tf.add_paragraph()
    p2.text = "Precision Agriculture Carbon Intelligence Platform"
    p2.font.name = 'Arial'
    p2.font.size = Pt(22)
    p2.font.color.rgb = COLOR_PRIMARY
    p2.space_before = Pt(10)
    
    # Subtitle Box
    sub_box = slide1.shapes.add_textbox(Inches(1.2), Inches(3.8), Inches(11.0), Inches(2.0))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    
    p_sub1 = tf_sub.paragraphs[0]
    p_sub1.text = "Evaluating Net Carbon Footprints and Soil Sequestration using Machine Learning"
    p_sub1.font.name = 'Arial'
    p_sub1.font.size = Pt(14)
    p_sub1.font.color.rgb = COLOR_WHITE
    p_sub1.space_after = Pt(25)
    
    p_sub2 = tf_sub.add_paragraph()
    p_sub2.text = "Project Assessment  •  Academic Review Presentation"
    p_sub2.font.name = 'Arial'
    p_sub2.font.size = Pt(12)
    p_sub2.font.color.rgb = COLOR_TEXT_LIGHT
    
    p_sub3 = tf_sub.add_paragraph()
    p_sub3.text = "Location: Karnataka State (31 Priority Districts Study)"
    p_sub3.font.name = 'Arial'
    p_sub3.font.size = Pt(11)
    p_sub3.font.color.rgb = COLOR_TEXT_LIGHT
    
    # Speaker Notes
    slide1.notes_slide.notes_text_frame.text = (
        "Welcome to the project presentation of CarbonIntel, a Precision Agriculture Carbon Intelligence Platform.\n\n"
        "Global agriculture faces a dual challenge: maximizing yields to feed a growing population while minimizing environmental impacts.\n"
        "CarbonIntel bridges this gap by turning complex soil chemistry, fertilizer logs, and weather data into actionable carbon footprint predictions and optimization recommendations.\n"
        "This presentation will walk through the system's motivation, experimental methodology, dataset details, ML performance, and frontend architecture."
    )

    # ==========================================
    # SLIDE 2: Project Motivation & Background
    # ==========================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide2, COLOR_BG)
    add_slide_header(slide2, "Project Motivation & Climatological Background")
    add_slide_footer(slide2, 2)
    
    # Left Card: The Problem
    card_left = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.1))
    card_left.fill.solid()
    card_left.fill.fore_color.rgb = COLOR_WHITE
    card_left.line.color.rgb = RGBColor(226, 232, 240)
    
    tf_left = card_left.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = Inches(0.3)
    tf_left.margin_top = Inches(0.3)
    tf_left.margin_right = Inches(0.3)
    
    p_l1 = tf_left.paragraphs[0]
    p_l1.text = "The Climate-Agriculture Paradox"
    p_l1.font.bold = True
    p_l1.font.size = Pt(18)
    p_l1.font.color.rgb = COLOR_TEXT_DARK
    p_l1.space_after = Pt(14)
    
    bullets_left = [
        "Global Warming Impact: Agriculture represents a significant share of global emissions, particularly nitrous oxide (N2O) from fertilizers and methane (CH4) from crop cultivation.",
        "Fertilizer Volatility: Over-application of chemical nitrogen fertilizers (Urea, Ammonium Nitrate) degrades soil microbes and creates strong greenhouse gas release.",
        "Sequestration Gap: Healthy Soil Organic Carbon (SOC) acts as an atmospheric carbon sink. However, standard testing lacks tools to link chemistry directly to footprint metrics."
    ]
    for b in bullets_left:
        p_b = tf_left.add_paragraph()
        p_b.text = b
        p_b.font.size = Pt(12)
        p_b.font.color.rgb = COLOR_TEXT_DARK
        p_b.space_before = Pt(8)
        
    # Right Card: The Solution
    card_right = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.1))
    card_right.fill.solid()
    card_right.fill.fore_color.rgb = COLOR_WHITE
    card_right.line.color.rgb = RGBColor(226, 232, 240)
    
    tf_right = card_right.text_frame
    tf_right.word_wrap = True
    tf_right.margin_left = Inches(0.3)
    tf_right.margin_top = Inches(0.3)
    tf_right.margin_right = Inches(0.3)
    
    p_r1 = tf_right.paragraphs[0]
    p_r1.text = "The CarbonIntel Approach"
    p_r1.font.bold = True
    p_r1.font.size = Pt(18)
    p_r1.font.color.rgb = COLOR_PRIMARY
    p_r1.space_after = Pt(14)
    
    bullets_right = [
        "Data-Driven Modeling: Merges laboratory soil card metrics (N, P, K, pH, SOC) with regional weather datasets to model agricultural emissions.",
        "Predictive Modeling: Utilizes machine learning models trained on hybrid datasets (real + statistical) to predict precise footprint ratings (kg CO2e/ha).",
        "Actionable Simulator: Provides farmers and agronomists with what-if interactive simulation tools to test chemical-to-organic fertilizer offsets in real-time."
    ]
    for b in bullets_right:
        p_b = tf_right.add_paragraph()
        p_b.text = b
        p_b.font.size = Pt(12)
        p_b.font.color.rgb = COLOR_TEXT_DARK
        p_b.space_before = Pt(8)

    # Speaker Notes
    slide2.notes_slide.notes_text_frame.text = (
        "This slide covers the background and core problem. Agriculture contributes heavily to emissions, yet growers lack simple tools to calculate and optimize their footprint.\n\n"
        "CarbonIntel resolves this paradox by translating the chemical metrics of a standard soil card—like Nitrogen, pH, and Carbon—coupled with temperature and rainfall, into actionable carbon estimates.\n"
        "This enables growers to see how management practices affect their score before committing resources."
    )

    # ==========================================
    # SLIDE 3: Problem Statement & Objectives
    # ==========================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3, COLOR_BG)
    add_slide_header(slide3, "Problem Statement & Research Objectives")
    add_slide_footer(slide3, 3)
    
    # 3 objective blocks horizontally
    obj_data = [
        ("Problem Statement", "Soil tests measure chemical nutrients but leave growers in the dark about greenhouse gas emissions. No accessible digital platform links local chemical soil profiles and live climatology to carbon footprint benchmarks (kg CO2e/ha).", COLOR_TEXT_DARK),
        ("Research Objective 1", "Develop and evaluate an optimized Machine Learning pipeline (using XGBoost and Random Forest) capable of predicting farm-level net carbon footprints from heterogeneous soil, crop, and weather factors.", COLOR_PRIMARY),
        ("Research Objective 2", "Implement an interactive, high-performance web platform combining satellite soil databases (ISRIC SoilGrids), NASA POWER meteorology, and lab directories to offer real-time optimization simulations.", COLOR_ACCENT)
    ]
    
    for i, (title, desc, color) in enumerate(obj_data):
        card = slide3.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i*3.9), Inches(1.8), Inches(3.7), Inches(4.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_WHITE
        card.line.color.rgb = RGBColor(226, 232, 240)
        
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.3)
        
        # Color bar at top of card
        bar = slide3.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8 + i*3.9), Inches(1.8), Inches(3.7), Inches(0.12)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = color
        p.space_after = Pt(12)
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12.5)
        p2.font.color.rgb = COLOR_TEXT_DARK
        p2.space_before = Pt(8)
        
    # Speaker Notes
    slide3.notes_slide.notes_text_frame.text = (
        "Here we define the core research problem and the primary objectives.\n\n"
        "Objective 1 is to train and validate a robust machine learning regressor that outperforms simple linear baselines in predicting carbon footprint.\n"
        "Objective 2 is to deploy this model inside a practical frontend tool that queries APIs like NASA and SoilGrids, helping users bypass manual weather log entries."
    )

    # ==========================================
    # SLIDE 4: System Architecture
    # ==========================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4, COLOR_BG)
    add_slide_header(slide4, "System Architecture & Integration Pipeline")
    add_slide_footer(slide4, 4)
    
    # Draw Architecture Blocks (Inputs, Services, ML, Outputs)
    blocks = [
        ("1. USER INPUTS", "District Selector\nMap Coordinate Pin\nManual Override", Inches(0.8), COLOR_TEXT_DARK),
        ("2. APIs & DIRECTORY", "NASA POWER API\nISRIC SoilGrids API\nAIKosh Lab Directory", Inches(3.8), COLOR_PRIMARY),
        ("3. ML INFERENCE", "FastAPI API endpoint\nXGBoost Pipeline\nStandardized Scale/Encode", Inches(6.8), COLOR_PRIMARY),
        ("4. RESULTS ENGINE", "Carbon Footprint\nSustainability rating\nWhat-If Simulator", Inches(9.8), COLOR_ACCENT)
    ]
    
    for title, desc, x, color in blocks:
        # Background card
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.2), Inches(2.7), Inches(3.6))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_WHITE
        card.line.color.rgb = RGBColor(226, 232, 240)
        
        # Banner inside card
        banner = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.2), Inches(2.7), Inches(0.6))
        banner.fill.solid()
        banner.fill.fore_color.rgb = color
        banner.line.fill.background()
        
        # Banner Title
        tf_b = banner.text_frame
        p_b = tf_b.paragraphs[0]
        p_b.text = title
        p_b.font.bold = True
        p_b.font.size = Pt(11)
        p_b.font.color.rgb = COLOR_WHITE
        p_b.alignment = PP_ALIGN.CENTER
        
        # Body text
        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_top = Inches(0.8)
        tf_c.margin_left = Inches(0.15)
        tf_c.margin_right = Inches(0.15)
        
        p_c = tf_c.paragraphs[0]
        p_c.text = desc
        p_c.font.size = Pt(13)
        p_c.font.color.rgb = COLOR_TEXT_DARK
        p_c.alignment = PP_ALIGN.CENTER
        p_c.line_spacing = 1.3
        
    # Draw simple connector arrows
    for i in range(3):
        arrow = slide4.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, Inches(3.55 + i*3.0), Inches(3.7), Inches(0.2), Inches(0.3)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = COLOR_TEXT_LIGHT
        arrow.line.fill.background()

    # Speaker Notes
    slide4.notes_slide.notes_text_frame.text = (
        "This block diagram displays the end-to-end architecture.\n\n"
        "First, a user inputs a location (via address, coordinate pin, or district select).\n"
        "Second, APIs fetch soil and weather values dynamically, resolving them into N, P, K, pH, temperature, and rain.\n"
        "Third, the React client passes this payload to the FastAPI backend, where the XGBoost model runs predictions.\n"
        "Finally, the results are rendered in the dashboard along with sustainability classifications and custom recommendation lists."
    )

    # ==========================================
    # SLIDE 5: Data Provenance & Preprocessing
    # ==========================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide5, COLOR_BG)
    add_slide_header(slide5, "Data Provenance & Cleaning Process")
    add_slide_footer(slide5, 5)
    
    # Left Content
    left_box = slide5.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.0))
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    
    p = tf_l.paragraphs[0]
    p.text = "Dataset Composition & Origin"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_TEXT_DARK
    p.space_after = Pt(12)
    
    bullets = [
        "Chitradurga Real Soil Cards: Sourced 17,785 real chemical soil samples from Challakere Taluk, providing empirical baseline metrics (avl_n, avl_p, avl_k, avl_oc, ph).",
        "Imputation & Cleanup: Missing chemical values were imputed using column-wise averages and clipped to realistic physical thresholds (e.g. pH strictly in [4.5, 8.5]).",
        "District Synthesis: Generated 2,000 synthesized records per district across 10 other priority Karnataka districts using Gaussian distributions based on official soil health card reports.",
        "External Mapped Data: Downloaded and structured agricultural dataset (originally crop recommendations) to expand baseline crop diversity."
    ]
    for b in bullets:
        p_b = tf_l.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(12.5)
        p_b.font.color.rgb = COLOR_TEXT_DARK
        p_b.space_before = Pt(8)
        
    # Right Table
    table_shape = slide5.shapes.add_table(6, 4, Inches(6.9), Inches(1.7), Inches(5.6), Inches(4.5))
    table = table_shape.table
    
    headers = ["Data Source", "Type", "Records", "Percentage"]
    for col_idx, h in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_HEADER_BG
        p_c = cell.text_frame.paragraphs[0]
        p_c.text = h
        p_c.font.bold = True
        p_c.font.size = Pt(11)
        p_c.font.color.rgb = COLOR_WHITE
        p_c.alignment = PP_ALIGN.CENTER
        
    data = [
        ["Chitradurga", "Real Soil Samples", "17,785", "47.2%"],
        ["Synthesized Districts", "Gaussian Mapped", "20,000", "52.8%"],
        ["Total Training", "Hybrid Compiled", "37,785", "100.0%"],
        ["SOC Boundary", "Clipped Range", "0.5% - 5.0%", "N/A"],
        ["pH Boundary", "Clipped Range", "4.5 - 8.5", "N/A"]
    ]
    
    for row_idx, r in enumerate(data):
        for col_idx, val in enumerate(r):
            cell = table.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            # Alternating background colors for rows
            if row_idx == 2:
                cell.fill.fore_color.rgb = RGBColor(209, 250, 229) # Light green for total
            elif row_idx % 2 == 0:
                cell.fill.fore_color.rgb = COLOR_WHITE
            else:
                cell.fill.fore_color.rgb = RGBColor(241, 245, 249)
            p_c = cell.text_frame.paragraphs[0]
            p_c.text = val
            p_c.font.size = Pt(10)
            p_c.font.color.rgb = COLOR_TEXT_DARK
            p_c.alignment = PP_ALIGN.CENTER
            
    # Speaker Notes
    slide5.notes_slide.notes_text_frame.text = (
        "This slide covers data collection and preprocessing.\n\n"
        "We used a hybrid dataset. 47.2% comes from real soil testing records in Chitradurga/Challakere (over 17,000 rows).\n"
        "To allow the model to generalize across Karnataka, we synthesized 2,000 records per district for 10 other districts using Gaussian distributions aligned with official soil health card reports.\n"
        "Values were clipped to strict biological limits to prevent the model from learning from anomalous outliers."
    )

    # ==========================================
    # SLIDE 6: Machine Learning Pipeline
    # ==========================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide6, COLOR_BG)
    add_slide_header(slide6, "Machine Learning Preprocessing & Pipeline")
    add_slide_footer(slide6, 6)
    
    # Draw vertical pipeline stages
    stages = [
        ("1. Input Features (11 features)", "Numerical (9): SOC, N, P, K, pH, Fertilizer Amount, Temp, Rain, Humidity.\nCategorical (2): Crop Type, Fertilizer Type."),
        ("2. ColumnTransformer Pipeline", "StandardScaler applied to scale all numerical ranges to zero-mean and unit-variance.\nOneHotEncoder (handle_unknown='ignore') applied to categorical variables."),
        ("3. Model Training & Selection", "Trained multiple estimators: Linear Regression (Baseline), Random Forest, XGBoost Regressor.\nHyperparameters tuned via GridSearch (learning_rate=0.07, max_depth=6, n_estimators=300).")
    ]
    
    for idx, (title, desc) in enumerate(stages):
        # Stage card
        card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.5 + idx*1.75), Inches(10.9), Inches(1.4))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_WHITE
        card.line.color.rgb = RGBColor(226, 232, 240)
        
        # Left color block
        block = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(1.5 + idx*1.75), Inches(0.3), Inches(1.4))
        block.fill.solid()
        block.fill.fore_color.rgb = COLOR_PRIMARY if idx != 2 else COLOR_ACCENT
        block.line.fill.background()
        
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.5)
        tf.margin_top = Inches(0.15)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXT_DARK
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_LIGHT
        p2.space_before = Pt(6)
        
        # Arrow connecting stages
        if idx < 2:
            arrow = slide6.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.5), Inches(2.95 + idx*1.75), Inches(0.2), Inches(0.25))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_TEXT_LIGHT
            arrow.line.fill.background()

    # Speaker Notes
    slide6.notes_slide.notes_text_frame.text = (
        "This slide visualizes the machine learning pipeline structure.\n\n"
        "We ingest 11 features: 9 numeric and 2 categorical.\n"
        "We use scikit-learn's ColumnTransformer to apply StandardScaler to numerical features, ensuring they are at the same scale, and OneHotEncoder to convert categorical variables into binary vectors.\n"
        "We then serialize this entire pipeline using joblib so that raw data can be passed directly to the predict API, avoiding training-serving skew."
    )

    # ==========================================
    # SLIDE 7: Model Evaluation & Results
    # ==========================================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide7, COLOR_BG)
    add_slide_header(slide7, "Experimental Results & Model Evaluation")
    add_slide_footer(slide7, 7)
    
    # Left explanation text
    text_box = slide7.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.0))
    tf_txt = text_box.text_frame
    tf_txt.word_wrap = True
    
    p_t = tf_txt.paragraphs[0]
    p_t.text = "Key Performance Findings"
    p_t.font.bold = True
    p_t.font.size = Pt(18)
    p_t.font.color.rgb = COLOR_TEXT_DARK
    p_t.space_after = Pt(14)
    
    bullets_t = [
        "Non-Linear superiority: Trees significantly outperform Linear Regression. The baseline Linear model achieved an R² of 0.9689, whereas XGBoost reached 0.9966.",
        "Why XGBoost Wins: The physical carbon emissions model features highly non-linear interactions, such as pH anomalies behaving quadratically, and multiplicative crop-fertilizer baselines.",
        "Error Reduction: Tuning XGBoost parameters (learning rate and depth) reduced Mean Absolute Error (MAE) from 22.43 kg/ha to 21.36 kg/ha, minimizing footprint prediction discrepancies."
    ]
    for b in bullets_t:
        p_b = tf_txt.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(12)
        p_b.font.color.rgb = COLOR_TEXT_DARK
        p_b.space_before = Pt(10)
        
    # Right Table: Model Comparison
    table_shape = slide7.shapes.add_table(5, 4, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.2))
    table_m = table_shape.table
    
    headers_m = ["Model Config", "MAE (kg/ha)", "RMSE (kg/ha)", "R-squared (R²)"]
    for col_idx, h in enumerate(headers_m):
        cell = table_m.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_HEADER_BG
        p_c = cell.text_frame.paragraphs[0]
        p_c.text = h
        p_c.font.bold = True
        p_c.font.size = Pt(11)
        p_c.font.color.rgb = COLOR_WHITE
        p_c.alignment = PP_ALIGN.CENTER
        
    model_rows = [
        ["Linear Regression", "59.58", "81.34", "0.9689"],
        ["Random Forest", "24.19", "31.26", "0.9954"],
        ["XGBoost (Default)", "22.43", "28.52", "0.9962"],
        ["XGBoost (Tuned)", "21.36", "27.06", "0.9966"]
    ]
    
    for row_idx, r in enumerate(model_rows):
        for col_idx, val in enumerate(r):
            cell = table_m.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            # Highlight Tuned XGBoost
            if row_idx == 3:
                cell.fill.fore_color.rgb = RGBColor(209, 250, 229) # Light Emerald Accent
            elif row_idx % 2 == 0:
                cell.fill.fore_color.rgb = COLOR_WHITE
            else:
                cell.fill.fore_color.rgb = RGBColor(241, 245, 249)
            p_c = cell.text_frame.paragraphs[0]
            p_c.text = val
            p_c.font.size = Pt(10.5)
            if row_idx == 3:
                p_c.font.bold = True
            p_c.font.color.rgb = COLOR_TEXT_DARK
            p_c.alignment = PP_ALIGN.CENTER

    # Speaker Notes
    slide7.notes_slide.notes_text_frame.text = (
        "Here are the core evaluation results.\n\n"
        "Linear Regression struggles to model complex physics. Trees do exceptionally well because they handle threshold cutoffs naturally.\n"
        "Our tuned XGBoost model achieves an R-squared of 99.66% and an MAE of 21.36 kg CO2e/ha.\n"
        "This small error margin is critical because a difference of 50 kg can decide whether a farm meets the 800 kg threshold to earn carbon credits."
    )

    # ==========================================
    # SLIDE 8: Core Platform Features
    # ==========================================
    slide8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide8, COLOR_BG)
    add_slide_header(slide8, "Core Features of the Web Application")
    add_slide_footer(slide8, 8)
    
    # 4 Cards for features
    features = [
        ("🌍 Live Geocoding Map", "Integrates interactive Leaflet maps and Open-Meteo search to pinpoint exact latitudes and longitudes, automatically mapping the region's default district profile.", Inches(0.8), COLOR_PRIMARY),
        ("🧪 Hybrid Soil Database", "Auto-queries coordinates against the global ISRIC SoilGrids API for SOC and pH. If local soil cards are available, users can manually override values.", Inches(3.7), COLOR_PRIMARY),
        ("🌦️ Climatology Fetcher", "Queries the NASA POWER Climatology API to retrieve 30-year historical weather data (annual precipitation, humidity, temperature), bypassing manual input logs.", Inches(6.6), COLOR_PRIMARY),
        ("🧬 Lab Proximity Finder", "Indexes 467 physical Indian Soil Testing Labs. Computes geodesic distances instantly using the Haversine formula to identify the nearest laboratory.", Inches(9.5), COLOR_ACCENT)
    ]
    
    for title, desc, x, color in features:
        card = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(2.7), Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_WHITE
        card.line.color.rgb = RGBColor(226, 232, 240)
        
        # Icon/Colored circle placeholder on card
        circle = slide8.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.0), Inches(2.1), Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(1.2)
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_TEXT_DARK
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(10)
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = COLOR_TEXT_LIGHT
        p2.alignment = PP_ALIGN.CENTER
        p2.line_spacing = 1.3

    # Speaker Notes
    slide8.notes_slide.notes_text_frame.text = (
        "Let's look at the web application capabilities.\n\n"
        "Instead of requiring manual entry for everything, CarbonIntel integrates dynamic geocoding, global soil databases, and climatology APIs to auto-populate features.\n"
        "Additionally, we compiled a directory of 467 Indian Soil Testing Laboratories. The client computes the Haversine distance, matching the user's location to the nearest physical lab to encourage empirical testing."
    )

    # ==========================================
    # SLIDE 9: Sustainability Rating & Carbon Credits
    # ==========================================
    slide9 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide9, COLOR_BG)
    add_slide_header(slide9, "Sustainability Ratings & Carbon Credits offset")
    add_slide_footer(slide9, 9)
    
    # Left Card: Sustainability Zones
    card_l = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.1))
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = COLOR_WHITE
    card_l.line.color.rgb = RGBColor(226, 232, 240)
    
    tf_l = card_l.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = Inches(0.3)
    tf_l.margin_top = Inches(0.3)
    tf_l.margin_right = Inches(0.3)
    
    p_l1 = tf_l.paragraphs[0]
    p_l1.text = "Sustainability Rating Categories"
    p_l1.font.bold = True
    p_l1.font.size = Pt(18)
    p_l1.font.color.rgb = COLOR_TEXT_DARK
    p_l1.space_after = Pt(14)
    
    zones = [
        ("High (Gold): < 400 kg CO2e/ha", "Assigned to highly efficient, low-emission practices (like organic vegetable/soybean farming with high SOC).", COLOR_PRIMARY),
        ("Medium (Silver): 400 - 1200 kg CO2e/ha", "Standard regional farming with moderate chemical fertilizer input. Suggests standard optimizations.", COLOR_ACCENT),
        ("Low (Bronze): > 1200 kg CO2e/ha", "Intensive chemical usage or highly anaerobic crops (flooded Rice). Requires immediate mitigation actions.", COLOR_TEXT_DARK)
    ]
    for title, desc, color in zones:
        p_t = tf_l.add_paragraph()
        p_t.text = title
        p_t.font.bold = True
        p_t.font.size = Pt(12)
        p_t.font.color.rgb = color
        p_t.space_before = Pt(8)
        
        p_d = tf_l.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = COLOR_TEXT_LIGHT
        p_d.space_after = Pt(6)
        
    # Right Card: Carbon Credits
    card_r = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.1))
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = COLOR_WHITE
    card_r.line.color.rgb = RGBColor(226, 232, 240)
    
    tf_r = card_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = Inches(0.3)
    tf_r.margin_top = Inches(0.3)
    tf_r.margin_right = Inches(0.3)
    
    p_r1 = tf_r.paragraphs[0]
    p_r1.text = "Carbon Credits Calculation"
    p_r1.font.bold = True
    p_r1.font.size = Pt(18)
    p_r1.font.color.rgb = COLOR_PRIMARY
    p_r1.space_after = Pt(14)
    
    # Formula Box
    f_box = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.2), Inches(2.2), Inches(5.0), Inches(1.0))
    f_box.fill.solid()
    f_box.fill.fore_color.rgb = RGBColor(241, 245, 249)
    f_box.line.color.rgb = RGBColor(203, 213, 225)
    tf_f = f_box.text_frame
    p_f = tf_f.paragraphs[0]
    p_f.text = "Credits (tCO2e) = Max(0, (800 - Footprint) / 1000)"
    p_f.font.name = 'Courier New'
    p_f.font.bold = True
    p_f.font.size = Pt(13.5)
    p_f.font.color.rgb = COLOR_TEXT_DARK
    p_f.alignment = PP_ALIGN.CENTER
    
    p_r2 = tf_r.add_paragraph()
    p_r2.text = "\nLogical Principles:"
    p_r2.font.bold = True
    p_r2.font.size = Pt(12)
    p_r2.font.color.rgb = COLOR_TEXT_DARK
    p_r2.space_before = Pt(25)
    
    credits_bullets = [
        "Benchmark Threshold: Uses a standard target ceiling of 800 kg CO2e/ha. Farms emitting below this benchmark generate offset credits.",
        "Credit Conversion: 1 Carbon Credit equals 1 metric tonne (1,000 kg) of CO2 equivalent offset.",
        "No Negative Penalty: High-emitting farms (> 800 kg) default to 0.00 credits, rather than suffering negative credit balances."
    ]
    for b in credits_bullets:
        p_b = tf_r.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(11)
        p_b.font.color.rgb = COLOR_TEXT_DARK
        p_b.space_before = Pt(8)

    # Speaker Notes
    slide9.notes_slide.notes_text_frame.text = (
        "This slide covers the translation of footprint metrics into economic incentives.\n\n"
        "The sustainability zones classify the farm output into Gold, Silver, and Bronze ratings to give clear environmental markers.\n"
        "The Carbon Credits formula rewards farmers who drop their footprint below the 800 kg benchmark. Every 1,000 kg offset represents one metric tonne of carbon offset, which translates directly to 1 carbon credit, incentivizing low-tillage, organic methods."
    )

    # ==========================================
    # SLIDE 10: Limitations & Future Scope
    # ==========================================
    slide10 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide10, COLOR_BG)
    add_slide_header(slide10, "Limitations & Future Scope")
    add_slide_footer(slide10, 10)
    
    # Left Card: Limitations
    card_l = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.1))
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = COLOR_WHITE
    card_l.line.color.rgb = RGBColor(226, 232, 240)
    
    tf_l = card_l.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = Inches(0.3)
    tf_l.margin_top = Inches(0.3)
    tf_l.margin_right = Inches(0.3)
    
    p_l1 = tf_l.paragraphs[0]
    p_l1.text = "Current Limitations"
    p_l1.font.bold = True
    p_l1.font.size = Pt(18)
    p_l1.font.color.rgb = COLOR_TEXT_DARK
    p_l1.space_after = Pt(14)
    
    l_bullets = [
        "Synthetic Target labels: The carbon footprint target variable is generated using standard emissions formulas rather than direct physical field gas measurements (due to lack of regional IoT carbon tracking sensors).",
        "Geographic Generalization: The real soil testing dataset represents Challakere Taluk (Chitradurga). Generalizing to other districts relies on global averages, which may introduce errors in highly unique local micro-climates."
    ]
    for b in l_bullets:
        p_b = tf_l.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(12)
        p_b.font.color.rgb = COLOR_TEXT_DARK
        p_b.space_before = Pt(8)
        
    # Right Card: Future Scope
    card_r = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.1))
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = COLOR_WHITE
    card_r.line.color.rgb = RGBColor(226, 232, 240)
    
    tf_r = card_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = Inches(0.3)
    tf_r.margin_top = Inches(0.3)
    tf_r.margin_right = Inches(0.3)
    
    p_r1 = tf_r.paragraphs[0]
    p_r1.text = "Future Directions"
    p_r1.font.bold = True
    p_r1.font.size = Pt(18)
    p_r1.font.color.rgb = COLOR_PRIMARY
    p_r1.space_after = Pt(14)
    
    f_bullets = [
        "IoT Integrations: Deploy real carbon dioxide (CO2) and nitrous oxide (N2O) sensors in test fields to collect real-world footprint data to retrain the ML model.",
        "Remote Sensing: Integrate Sentinel-2 / Landsat satellite imagery to capture NDVI and biomass data to estimate SOC variations dynamically over time.",
        "Statewide Expansion: Connect to other Indian state soil databases (like Tamil Nadu or Maharashtra) to broaden the scope outside Karnataka."
    ]
    for b in f_bullets:
        p_b = tf_r.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(12)
        p_b.font.color.rgb = COLOR_TEXT_DARK
        p_b.space_before = Pt(8)

    # Speaker Notes
    slide10.notes_slide.notes_text_frame.text = (
        "Every project has limitations. The primary limitation is the lack of physical field-level carbon measurements, forcing us to rely on IPCC formulas to model the target variable.\n\n"
        "In the future, we plan to validate this model using physical carbon-chamber gas measurements.\n"
        "We also plan to integrate satellite Sentinel-2 data to measure crop biomass, which directly correlates with carbon sequestration indices, moving towards a fully automated tracking system."
    )

    # ==========================================
    # SLIDE 11: References & Acknowledgements
    # ==========================================
    slide11 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide11, COLOR_BG)
    add_slide_header(slide11, "References & Acknowledgements")
    add_slide_footer(slide11, 11)
    
    # Text box for references
    ref_box = slide11.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.0))
    tf_ref = ref_box.text_frame
    tf_ref.word_wrap = True
    
    p_ref = tf_ref.paragraphs[0]
    p_ref.text = "Key Data & Literature References"
    p_ref.font.bold = True
    p_ref.font.size = Pt(18)
    p_ref.font.color.rgb = COLOR_TEXT_DARK
    p_ref.space_after = Pt(14)
    
    references = [
        "1. Intergovernmental Panel on Climate Change (IPCC) - '2019 Refinement to the 2006 IPCC Guidelines for National Greenhouse Gas Inventories' (Agricultural Soils chapter).",
        "2. Indian Soil Health Card (SHC) Portal - Ministry of Agriculture & Farmers Welfare (Karnataka State soil records).",
        "3. NASA POWER Project (Prediction Of Worldwide Energy Resources) - Climatological daily and monthly planetary atmospheric profiles.",
        "4. ISRIC World Soil Information - 'SoilGrids 2.0: Global soil property predictions map at 250 m resolution'.",
        "5. AIKosh Soil Testing Lab Directory database (Geographic listings of active agricultural testing centers in India)."
    ]
    for r in references:
        p_r = tf_ref.add_paragraph()
        p_r.text = r
        p_r.font.size = Pt(12)
        p_r.font.color.rgb = COLOR_TEXT_DARK
        p_r.space_before = Pt(8)
        
    p_ack = tf_ref.add_paragraph()
    p_ack.text = "\nAcknowledgements"
    p_ack.font.bold = True
    p_ack.font.size = Pt(16)
    p_ack.font.color.rgb = COLOR_PRIMARY
    p_ack.space_before = Pt(20)
    
    p_ack_text = tf_ref.add_paragraph()
    p_ack_text.text = "We express gratitude to the open-science initiatives of NASA, ISRIC, Open-Meteo, and the open-data releases from the Department of Agriculture, Government of Karnataka, for providing the scientific baseline datasets that made this validation study possible."
    p_ack_text.font.size = Pt(11.5)
    p_ack_text.font.color.rgb = COLOR_TEXT_LIGHT
    p_ack_text.space_before = Pt(6)

    # Speaker Notes
    slide11.notes_slide.notes_text_frame.text = (
        "We conclude by listing our scientific and data references, including the IPCC guidelines, NASA POWER, ISRIC SoilGrids, and the Indian Soil Health Card database.\n\n"
        "We thank our mentors, evaluators, and the open data initiatives that made this study possible.\n"
        "Thank you for your time. I am now open to any questions regarding the ML modeling, API ingestion layers, or platform integration."
    )

    # Save presentation
    output_dir = 'artifacts'
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, 'carbonintel_project_presentation.pptx')
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == '__main__':
    create_presentation()
